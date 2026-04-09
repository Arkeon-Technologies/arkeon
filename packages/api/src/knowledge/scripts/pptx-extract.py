#!/usr/bin/env python3
"""
Extract text, tables, notes, and images from a PPTX file.
Outputs structured JSON to stdout.

Usage: python3 pptx-extract.py <pptx_path> --image-dir <dir>
"""

import argparse
import json
import os
import sys

from pptx import Presentation
from pptx.enum.shapes import MSO_SHAPE_TYPE


def resize_image_if_needed(image_bytes, content_type, max_dim=2048):
    """Cap image at max_dim px on longest side using Pillow."""
    try:
        from PIL import Image
        import io

        img = Image.open(io.BytesIO(image_bytes))
        w, h = img.size
        if max(w, h) <= max_dim:
            return image_bytes, content_type

        scale = max_dim / max(w, h)
        new_w, new_h = int(w * scale), int(h * scale)
        img = img.resize((new_w, new_h), Image.LANCZOS)

        buf = io.BytesIO()
        out_format = "JPEG"
        out_ct = "image/jpeg"
        if content_type == "image/png":
            out_format = "PNG"
            out_ct = "image/png"
        img.save(buf, format=out_format)
        return buf.getvalue(), out_ct
    except Exception:
        return image_bytes, content_type


def extract_slide(slide, slide_number, image_dir):
    """Extract text, tables, notes, and images from a single slide."""
    text_parts = []
    tables = []
    images = []
    img_counter = 0

    for shape in slide.shapes:
        # Text frames
        if shape.has_text_frame:
            frame_text = shape.text_frame.text.strip()
            if frame_text:
                text_parts.append(frame_text)

        # Tables
        if shape.has_table:
            table = shape.table
            rows = []
            for row in table.rows:
                cells = [cell.text.strip() for cell in row.cells]
                rows.append(cells)
            if rows:
                tables.append({"rows": rows})

        # Images
        if shape.shape_type == MSO_SHAPE_TYPE.PICTURE:
            try:
                image = shape.image
                blob = image.blob
                ct = image.content_type or "image/png"

                ext = "png"
                if "jpeg" in ct or "jpg" in ct:
                    ext = "jpg"
                elif "gif" in ct:
                    ext = "gif"
                elif "bmp" in ct:
                    ext = "bmp"

                blob, ct = resize_image_if_needed(blob, ct)

                filename = f"slide_{slide_number}_img_{img_counter}.{ext}"
                path = os.path.join(image_dir, filename)
                with open(path, "wb") as f:
                    f.write(blob)

                images.append({
                    "filename": filename,
                    "content_type": ct,
                    "path": path,
                })
                img_counter += 1
            except Exception:
                pass

    # Speaker notes
    notes = ""
    try:
        if slide.has_notes_slide and slide.notes_slide.notes_text_frame:
            notes = slide.notes_slide.notes_text_frame.text.strip()
    except Exception:
        pass

    return {
        "slide_number": slide_number,
        "text": "\n\n".join(text_parts),
        "notes": notes,
        "tables": tables,
        "images": images,
    }


def main():
    parser = argparse.ArgumentParser(description="Extract content from PPTX")
    parser.add_argument("pptx_path", help="Path to .pptx file")
    parser.add_argument("--image-dir", required=True, help="Directory for extracted images")
    args = parser.parse_args()

    if not os.path.exists(args.pptx_path):
        print(f"File not found: {args.pptx_path}", file=sys.stderr)
        sys.exit(1)

    os.makedirs(args.image_dir, exist_ok=True)

    prs = Presentation(args.pptx_path)
    slides = []

    for i, slide in enumerate(prs.slides, start=1):
        slides.append(extract_slide(slide, i, args.image_dir))

    result = {
        "slide_count": len(slides),
        "slides": slides,
    }

    json.dump(result, sys.stdout)


if __name__ == "__main__":
    main()
